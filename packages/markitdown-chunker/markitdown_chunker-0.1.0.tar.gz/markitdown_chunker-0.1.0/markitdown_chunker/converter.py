"""
Document to Markdown converter using markitdown package with image extraction.
"""

import os
import shutil
import base64
import re
from pathlib import Path
from typing import Optional, List, Dict
from markitdown import MarkItDown


class MarkdownConverter:
    """
    Converts various document formats to Markdown using markitdown.
    
    Supported formats: pdf, docx, doc, pptx, ppt, xlsx, xls, 
                      html, htm, txt, md, rtf, odt, ods, odp
    """
    
    SUPPORTED_FORMATS = {
        '.pdf', '.docx', '.doc', '.pptx', '.ppt', 
        '.xlsx', '.xls', '.html', '.htm', '.txt',
        '.md', '.rtf', '.odt', '.ods', '.odp'
    }
    
    def __init__(self):
        """Initialize the MarkdownConverter."""
        self.markitdown = MarkItDown()
    
    def is_supported(self, file_path: str) -> bool:
        """
        Check if the file format is supported.
        
        Args:
            file_path: Path to the file
            
        Returns:
            True if supported, False otherwise
        """
        ext = Path(file_path).suffix.lower()
        return ext in self.SUPPORTED_FORMATS
    
    def convert(
        self, 
        file_path: str, 
        output_dir: str,
        save_images: bool = True
    ) -> dict:
        """
        Convert a document to Markdown format.
        
        Args:
            file_path: Path to the input file
            output_dir: Directory to save the output
            save_images: Whether to save extracted images
            
        Returns:
            Dictionary containing:
                - markdown_path: Path to the saved markdown file
                - images_dir: Path to the images directory (if images were saved)
                - markdown_content: The markdown content as string
                
        Raises:
            FileNotFoundError: If the input file doesn't exist
            ValueError: If the file format is not supported
        """
        # Validate input
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        if not self.is_supported(file_path):
            raise ValueError(
                f"Unsupported file format. Supported formats: {self.SUPPORTED_FORMATS}"
            )
        
        # Create output directory
        output_path = Path(output_dir)
        output_path.mkdir(parents=True, exist_ok=True)
        
        # Convert to markdown
        result = self.markitdown.convert(file_path)
        markdown_content = result.text_content
        
        # Generate output filename
        input_filename = Path(file_path).stem
        markdown_filename = f"{input_filename}.md"
        markdown_path = output_path / markdown_filename
        
        # Handle images if present
        images_dir = None
        extracted_images = []
        
        if save_images:
            images_dir = output_path / "images"
            images_dir.mkdir(exist_ok=True)
            
            # Extract images based on file type
            extracted_images = self._extract_images(file_path, images_dir)
            
            # Update markdown to reference extracted images
            if extracted_images:
                markdown_content = self._update_image_references(
                    markdown_content, 
                    extracted_images,
                    images_dir,
                    output_path
                )
        
        # Save markdown file (after image reference updates)
        with open(markdown_path, 'w', encoding='utf-8') as f:
            f.write(markdown_content)
        
        return {
            "markdown_path": str(markdown_path),
            "images_dir": str(images_dir) if images_dir else None,
            "markdown_content": markdown_content,
            "source_file": file_path,
            "output_dir": output_dir,
            "extracted_images": extracted_images
        }
    
    def convert_batch(
        self,
        file_paths: List[str],
        output_dir: str,
        save_images: bool = True
    ) -> List[dict]:
        """
        Convert multiple documents to Markdown.
        
        Args:
            file_paths: List of paths to input files
            output_dir: Directory to save the outputs
            save_images: Whether to save extracted images
            
        Returns:
            List of result dictionaries from convert()
        """
        results = []
        for file_path in file_paths:
            try:
                result = self.convert(file_path, output_dir, save_images)
                results.append(result)
            except Exception as e:
                results.append({
                    "error": str(e),
                    "source_file": file_path
                })
        return results
    
    def _extract_images(self, file_path: str, images_dir: Path) -> List[Dict[str, str]]:
        """
        Extract images from documents.
        
        Args:
            file_path: Path to the document
            images_dir: Directory to save extracted images
            
        Returns:
            List of dictionaries with image info (path, name, etc.)
        """
        ext = Path(file_path).suffix.lower()
        extracted = []
        
        try:
            # Extract images from PDF
            if ext == '.pdf':
                extracted = self._extract_pdf_images(file_path, images_dir)
            
            # Extract images from DOCX
            elif ext in ['.docx', '.doc']:
                extracted = self._extract_docx_images(file_path, images_dir)
            
            # Extract images from PPTX
            elif ext in ['.pptx', '.ppt']:
                extracted = self._extract_pptx_images(file_path, images_dir)
                
        except Exception as e:
            print(f"Warning: Could not extract images from {file_path}: {e}")
        
        return extracted
    
    def _extract_pdf_images(self, pdf_path: str, images_dir: Path) -> List[Dict[str, str]]:
        """Extract images from PDF using PyMuPDF if available."""
        extracted = []
        
        try:
            import fitz  # PyMuPDF
            
            pdf_document = fitz.open(pdf_path)
            
            for page_num in range(len(pdf_document)):
                page = pdf_document[page_num]
                image_list = page.get_images()
                
                for img_index, img in enumerate(image_list):
                    xref = img[0]
                    base_image = pdf_document.extract_image(xref)
                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]
                    
                    # Save image
                    image_filename = f"page{page_num + 1}_img{img_index + 1}.{image_ext}"
                    image_path = images_dir / image_filename
                    
                    with open(image_path, 'wb') as img_file:
                        img_file.write(image_bytes)
                    
                    extracted.append({
                        "filename": image_filename,
                        "path": str(image_path),
                        "page": page_num + 1,
                        "index": img_index + 1
                    })
            
            pdf_document.close()
            
        except ImportError:
            print("Info: Install 'pymupdf' for PDF image extraction: pip install pymupdf")
        except Exception as e:
            print(f"Warning: Error extracting PDF images: {e}")
        
        return extracted
    
    def _extract_docx_images(self, docx_path: str, images_dir: Path) -> List[Dict[str, str]]:
        """Extract images from DOCX using python-docx if available."""
        extracted = []
        
        try:
            from docx import Document
            from docx.opc.constants import RELATIONSHIP_TYPE as RT
            
            doc = Document(docx_path)
            
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    image_part = rel.target_part
                    image_bytes = image_part.blob
                    
                    # Determine extension from content type
                    content_type = image_part.content_type
                    ext_map = {
                        'image/jpeg': 'jpg',
                        'image/png': 'png',
                        'image/gif': 'gif',
                        'image/bmp': 'bmp',
                        'image/tiff': 'tiff'
                    }
                    ext = ext_map.get(content_type, 'jpg')
                    
                    # Save image
                    image_filename = f"docx_img{len(extracted) + 1}.{ext}"
                    image_path = images_dir / image_filename
                    
                    with open(image_path, 'wb') as img_file:
                        img_file.write(image_bytes)
                    
                    extracted.append({
                        "filename": image_filename,
                        "path": str(image_path),
                        "index": len(extracted) + 1
                    })
            
        except ImportError:
            print("Info: Install 'python-docx' for DOCX image extraction: pip install python-docx")
        except Exception as e:
            print(f"Warning: Error extracting DOCX images: {e}")
        
        return extracted
    
    def _extract_pptx_images(self, pptx_path: str, images_dir: Path) -> List[Dict[str, str]]:
        """Extract images from PPTX using python-pptx if available."""
        extracted = []
        
        try:
            from pptx import Presentation
            
            prs = Presentation(pptx_path)
            
            for slide_num, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if hasattr(shape, "image"):
                        image = shape.image
                        image_bytes = image.blob
                        
                        # Get extension
                        ext = image.ext or 'jpg'
                        
                        # Save image
                        image_filename = f"slide{slide_num + 1}_img{len(extracted) + 1}.{ext}"
                        image_path = images_dir / image_filename
                        
                        with open(image_path, 'wb') as img_file:
                            img_file.write(image_bytes)
                        
                        extracted.append({
                            "filename": image_filename,
                            "path": str(image_path),
                            "slide": slide_num + 1,
                            "index": len(extracted) + 1
                        })
            
        except ImportError:
            print("Info: Install 'python-pptx' for PPTX image extraction: pip install python-pptx")
        except Exception as e:
            print(f"Warning: Error extracting PPTX images: {e}")
        
        return extracted
    
    def _update_image_references(
        self, 
        markdown_content: str, 
        extracted_images: List[Dict[str, str]],
        images_dir: Path,
        output_dir: Path
    ) -> str:
        """
        Add image references to markdown content.
        
        Args:
            markdown_content: Original markdown text
            extracted_images: List of extracted image info
            images_dir: Directory where images are saved
            output_dir: Base output directory
            
        Returns:
            Updated markdown content with image references
        """
        if not extracted_images:
            return markdown_content
        
        # Add image section at the end
        image_section = "\n\n## Extracted Images\n\n"
        
        for img in extracted_images:
            # Create relative path from markdown file to image
            rel_path = f"images/{img['filename']}"
            
            # Add markdown image reference
            alt_text = f"Image {img.get('index', '')}"
            if 'page' in img:
                alt_text = f"Page {img['page']} Image {img['index']}"
            elif 'slide' in img:
                alt_text = f"Slide {img['slide']} Image {img['index']}"
            
            image_section += f"![{alt_text}]({rel_path})\n\n"
        
        return markdown_content + image_section

