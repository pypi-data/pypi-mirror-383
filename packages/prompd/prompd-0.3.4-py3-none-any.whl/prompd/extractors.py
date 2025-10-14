"""
Binary file extraction utilities for the Prompd compiler.

Supports extracting text content from various binary formats:
- Excel (.xlsx, .xls) → CSV/tabular data
- Word (.docx) → plain text
- PDF (.pdf) → OCR text extraction
- PowerPoint (.pptx) → slide text
- Images (.png, .jpg, .gif) → alt-text/descriptions
- Config files (.json, .yaml, .env) → structured data
"""

import json
import csv
import io
from pathlib import Path
from typing import Union, Dict, Any, Optional, List
from abc import ABC, abstractmethod

try:
    import openpyxl
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import PyPDF2
    PYPDF2_AVAILABLE = True
except ImportError:
    PYPDF2_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False

import yaml


class FileExtractor(ABC):
    """Abstract base class for file extractors."""
    
    @abstractmethod
    def can_extract(self, file_path: Path) -> bool:
        """Check if this extractor can handle the given file."""
        pass
    
    @abstractmethod
    def extract(self, file_path: Path) -> str:
        """Extract text content from the file."""
        pass
    
    @abstractmethod
    def get_supported_extensions(self) -> List[str]:
        """Get list of supported file extensions."""
        pass


class ExcelExtractor(FileExtractor):
    """Extract data from Excel files as CSV format."""
    
    def can_extract(self, file_path: Path) -> bool:
        """Check if file is an Excel file."""
        return OPENPYXL_AVAILABLE and file_path.suffix.lower() in ['.xlsx', '.xlsm']
    
    def extract(self, file_path: Path) -> str:
        """Extract Excel data as CSV format."""
        if not OPENPYXL_AVAILABLE:
            return f"# Excel extraction not available - install openpyxl\n# File: {file_path}"
        
        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            output = []
            
            for sheet_name in workbook.sheetnames:
                worksheet = workbook[sheet_name]
                
                if len(workbook.sheetnames) > 1:
                    output.append(f"# Sheet: {sheet_name}")
                    output.append("")
                
                # Convert sheet to CSV format
                csv_buffer = io.StringIO()
                csv_writer = csv.writer(csv_buffer)
                
                for row in worksheet.iter_rows(values_only=True):
                    # Skip completely empty rows
                    if any(cell is not None for cell in row):
                        # Convert None values to empty strings
                        clean_row = [str(cell) if cell is not None else "" for cell in row]
                        csv_writer.writerow(clean_row)
                
                output.append(csv_buffer.getvalue().strip())
                output.append("")
            
            return "\n".join(output).strip()
            
        except Exception as e:
            return f"# Error extracting Excel file: {e}\n# File: {file_path}"
    
    def get_supported_extensions(self) -> List[str]:
        return ['.xlsx', '.xlsm'] if OPENPYXL_AVAILABLE else []


class WordExtractor(FileExtractor):
    """Extract text from Word documents."""
    
    def can_extract(self, file_path: Path) -> bool:
        """Check if file is a Word document."""
        return DOCX_AVAILABLE and file_path.suffix.lower() == '.docx'
    
    def extract(self, file_path: Path) -> str:
        """Extract text from Word document."""
        if not DOCX_AVAILABLE:
            return f"# Word extraction not available - install python-docx\n# File: {file_path}"
        
        try:
            document = Document(file_path)
            
            paragraphs = []
            for paragraph in document.paragraphs:
                text = paragraph.text.strip()
                if text:  # Skip empty paragraphs
                    paragraphs.append(text)
            
            # Extract table content
            tables_content = []
            for table in document.tables:
                table_data = []
                for row in table.rows:
                    row_data = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        row_data.append(cell_text)
                    if any(cell for cell in row_data):  # Skip empty rows
                        table_data.append(" | ".join(row_data))
                
                if table_data:
                    tables_content.append("\n".join(table_data))
            
            content = []
            if paragraphs:
                content.extend(paragraphs)
            
            if tables_content:
                if content:
                    content.append("")  # Add separator
                content.append("# Tables")
                content.extend(tables_content)
            
            return "\n".join(content) if content else f"# No text content found in {file_path}"
            
        except Exception as e:
            return f"# Error extracting Word document: {e}\n# File: {file_path}"
    
    def get_supported_extensions(self) -> List[str]:
        return ['.docx'] if DOCX_AVAILABLE else []


class PDFExtractor(FileExtractor):
    """Extract text from PDF files."""
    
    def can_extract(self, file_path: Path) -> bool:
        """Check if file is a PDF."""
        return PYPDF2_AVAILABLE and file_path.suffix.lower() == '.pdf'
    
    def extract(self, file_path: Path) -> str:
        """Extract text from PDF."""
        if not PYPDF2_AVAILABLE:
            return f"# PDF extraction not available - install PyPDF2\n# File: {file_path}"
        
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                text_content = []
                for page_num, page in enumerate(pdf_reader.pages, 1):
                    try:
                        page_text = page.extract_text().strip()
                        if page_text:
                            if len(pdf_reader.pages) > 1:
                                text_content.append(f"# Page {page_num}")
                            text_content.append(page_text)
                    except Exception as e:
                        text_content.append(f"# Error extracting page {page_num}: {e}")
                
                return "\n\n".join(text_content) if text_content else f"# No text content found in {file_path}"
                
        except Exception as e:
            return f"# Error extracting PDF: {e}\n# File: {file_path}"
    
    def get_supported_extensions(self) -> List[str]:
        return ['.pdf'] if PYPDF2_AVAILABLE else []


class PowerPointExtractor(FileExtractor):
    """Extract text from PowerPoint presentations."""
    
    def can_extract(self, file_path: Path) -> bool:
        """Check if file is a PowerPoint presentation."""
        return PPTX_AVAILABLE and file_path.suffix.lower() in ['.pptx', '.pptm']
    
    def extract(self, file_path: Path) -> str:
        """Extract text from PowerPoint presentation."""
        if not PPTX_AVAILABLE:
            return f"# PowerPoint extraction not available - install python-pptx\n# File: {file_path}"
        
        try:
            presentation = Presentation(file_path)
            slides_content = []
            
            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_text = []
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    slides_content.append(f"# Slide {slide_num}")
                    slides_content.extend(slide_text)
                    slides_content.append("")  # Add spacing between slides
            
            return "\n".join(slides_content).strip() if slides_content else f"# No text content found in {file_path}"
            
        except Exception as e:
            return f"# Error extracting PowerPoint: {e}\n# File: {file_path}"
    
    def get_supported_extensions(self) -> List[str]:
        return ['.pptx', '.pptm'] if PPTX_AVAILABLE else []


class ImageExtractor(FileExtractor):
    """Extract metadata and descriptions from images."""
    
    def can_extract(self, file_path: Path) -> bool:
        """Check if file is an image."""
        return PIL_AVAILABLE and file_path.suffix.lower() in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp']
    
    def extract(self, file_path: Path) -> str:
        """Extract image metadata and description."""
        if not PIL_AVAILABLE:
            return f"# Image extraction not available - install Pillow\n# File: {file_path}"
        
        try:
            with Image.open(file_path) as image:
                info = []
                info.append(f"# Image: {file_path.name}")
                info.append(f"Format: {image.format}")
                info.append(f"Size: {image.size[0]}x{image.size[1]} pixels")
                info.append(f"Mode: {image.mode}")
                
                # Extract EXIF data if available
                if hasattr(image, '_getexif') and image._getexif():
                    info.append("# EXIF Data:")
                    exif = image._getexif()
                    for key, value in exif.items():
                        if isinstance(value, str) and len(value) < 100:  # Only include short string values
                            info.append(f"{key}: {value}")
                
                # TODO: In the future, integrate with vision API for image descriptions
                info.append("# Description: [Image content analysis not yet implemented]")
                
                return "\n".join(info)
                
        except Exception as e:
            return f"# Error extracting image metadata: {e}\n# File: {file_path}"
    
    def get_supported_extensions(self) -> List[str]:
        return ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp'] if PIL_AVAILABLE else []


class JSONExtractor(FileExtractor):
    """Extract structured data from JSON files."""
    
    def can_extract(self, file_path: Path) -> bool:
        """Check if file is JSON."""
        return file_path.suffix.lower() == '.json'
    
    def extract(self, file_path: Path) -> str:
        """Extract and format JSON data."""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                data = json.load(file)
                
                # Pretty format the JSON
                formatted = json.dumps(data, indent=2, ensure_ascii=False)
                return f"# JSON Data from {file_path.name}\n\n```json\n{formatted}\n```"
                
        except Exception as e:
            return f"# Error extracting JSON: {e}\n# File: {file_path}"
    
    def get_supported_extensions(self) -> List[str]:
        return ['.json']


class YAMLExtractor(FileExtractor):
    """Extract structured data from YAML files."""
    
    def can_extract(self, file_path: Path) -> bool:
        """Check if file is YAML."""
        return file_path.suffix.lower() in ['.yaml', '.yml']
    
    def extract(self, file_path: Path) -> str:
        """Extract and format YAML data."""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                data = yaml.safe_load(file)
                
                # Format as YAML
                formatted = yaml.dump(data, default_flow_style=False, allow_unicode=True)
                return f"# YAML Data from {file_path.name}\n\n```yaml\n{formatted}\n```"
                
        except Exception as e:
            return f"# Error extracting YAML: {e}\n# File: {file_path}"
    
    def get_supported_extensions(self) -> List[str]:
        return ['.yaml', '.yml']


class TextExtractor(FileExtractor):
    """Extract content from plain text files."""
    
    def can_extract(self, file_path: Path) -> bool:
        """Check if file is a text file."""
        return file_path.suffix.lower() in ['.txt', '.md', '.py', '.js', '.ts', '.sql', '.css', '.html', '.xml', '.env', '.csv']
    
    def extract(self, file_path: Path) -> str:
        """Extract text content."""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read().strip()
                
                # Add syntax highlighting hint for code files
                ext = file_path.suffix.lower()
                if ext in ['.py', '.js', '.ts', '.sql', '.css', '.html', '.xml', '.csv']:
                    lang_map = {
                        '.py': 'python',
                        '.js': 'javascript', 
                        '.ts': 'typescript',
                        '.sql': 'sql',
                        '.css': 'css',
                        '.html': 'html',
                        '.xml': 'xml',
                        '.csv': 'csv'
                    }
                    lang = lang_map.get(ext, '')
                    return f"# {file_path.name}\n\n```{lang}\n{content}\n```"
                else:
                    return f"# {file_path.name}\n\n{content}"
                
        except Exception as e:
            return f"# Error extracting text file: {e}\n# File: {file_path}"
    
    def get_supported_extensions(self) -> List[str]:
        return ['.txt', '.md', '.py', '.js', '.ts', '.sql', '.css', '.html', '.xml', '.env', '.csv']


class BinaryExtractionEngine:
    """Main engine for extracting content from binary files."""
    
    def __init__(self):
        self.extractors = [
            ExcelExtractor(),
            WordExtractor(), 
            PDFExtractor(),
            PowerPointExtractor(),
            ImageExtractor(),
            JSONExtractor(),
            YAMLExtractor(),
            TextExtractor(),  # Keep as fallback
        ]
    
    def can_extract(self, file_path: Union[str, Path]) -> bool:
        """Check if any extractor can handle this file."""
        file_path = Path(file_path)
        return any(extractor.can_extract(file_path) for extractor in self.extractors)
    
    def extract(self, file_path: Union[str, Path]) -> str:
        """Extract content from the file using the appropriate extractor."""
        file_path = Path(file_path)
        
        if not file_path.exists():
            return f"# File not found: {file_path}"
        
        # Find the first extractor that can handle this file
        for extractor in self.extractors:
            if extractor.can_extract(file_path):
                return extractor.extract(file_path)
        
        return f"# No extractor available for file type: {file_path.suffix}\n# File: {file_path}"
    
    def get_supported_extensions(self) -> Dict[str, List[str]]:
        """Get all supported extensions organized by extractor type."""
        supported = {}
        for extractor in self.extractors:
            extractor_name = extractor.__class__.__name__
            supported[extractor_name] = extractor.get_supported_extensions()
        return supported
    
    def extract_multiple(self, file_paths: List[Union[str, Path]]) -> Dict[str, str]:
        """Extract content from multiple files."""
        results = {}
        for file_path in file_paths:
            path_obj = Path(file_path)
            results[str(path_obj)] = self.extract(path_obj)
        return results


# Global instance for easy access
binary_extractor = BinaryExtractionEngine()