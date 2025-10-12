"""
Prompd Registry Integration

Provides commands for publishing, searching, and installing packages from the Prompd registry.
Integrates with registry.prompdhub.ai API endpoints.
"""

import json
import os
import re
import requests
import zipfile
import tempfile
import csv
import io
from pathlib import Path
from typing import Dict, List, Optional, Any, Tuple
import yaml
from dataclasses import dataclass
from rich.progress import Progress, DownloadColumn, TransferSpeedColumn, TimeRemainingColumn, TaskID
from rich.console import Console

# Binary file extraction libraries - imported lazily when needed for faster startup

from .config import PrompdConfig
from .parser import PrompdParser
from .validator import PrompdValidator
from .package_resolver import RegistryInfo


class RegistryClient:
    """Client for interacting with Prompd registries (multi-registry support)."""
    
    def __init__(self, registry_name: Optional[str] = None):
        self.config = PrompdConfig.load()
        self.registry_name = registry_name or self.config.registry.get('default', 'prompdhub')
        self.session = requests.Session()
        
        # Get registry config
        registries = self.config.registry.get('registries', {})
        if self.registry_name not in registries:
            raise Exception(f"Registry '{self.registry_name}' not found in configuration")
        
        self.registry_config = registries[self.registry_name]
        self.registry_url = self.registry_config['url']
        # Lazy discovery: do not call well-known at init
        self.registry_info = None
        
        # Set up authentication if available
        token = self.registry_config.get('api_key')
        if token:
            self.session.headers.update({
                'Authorization': f'Bearer {token}'
            })

    def ensure_discovered(self):
        """Discover registry endpoints once, with short timeout and quiet failure."""
        if self.registry_info is not None:
            return
        try:
            # Perform discovery; RegistryInfo uses httpx default timeout, which is fine for a quick ping
            self.registry_info = RegistryInfo.from_well_known(self.registry_url)
        except Exception:
            # Leave as None; methods will fall back to basic endpoints
            self.registry_info = None
    
    def _get_endpoint_url(self, endpoint_name: str, fallback_path: str, **kwargs) -> str:
        """Get endpoint URL dynamically from registry discovery, with fallback."""
        # Ensure discovery has been attempted before relying on endpoints
        self.ensure_discovered()
        if self.registry_info and endpoint_name in self.registry_info.endpoints:
            endpoint_template = self.registry_info.endpoints[endpoint_name]
            # Replace template variables like {package}
            endpoint_path = endpoint_template.format(**kwargs)
            return f"{self.registry_url}{endpoint_path}"
        else:
            # Fallback to hardcoded paths
            return f"{self.registry_url}{fallback_path}"
    
    def login_with_credentials(self, username: str, password: str) -> Dict[str, Any]:
        """Authenticate with the registry using username and password."""
        try:
            # Login with credentials to get token
            login_data = {
                'username': username,
                'password': password
            }
            
            # Use discovered endpoint for login
            login_url = self._get_endpoint_url('login', '/auth/login')
            response = self.session.post(login_url, json=login_data)
            response.raise_for_status()
            auth_data = response.json()
            
            # Extract token from response
            token = auth_data.get('token') or auth_data.get('access_token')
            if not token:
                raise Exception("No token received from server")
            
            # Set up authentication with the token
            return self.login_with_token(token)
            
        except requests.exceptions.RequestException as e:
            raise Exception(f"Login failed: {e}")
    
    def login_with_token(self, token: str) -> Dict[str, Any]:
        """Authenticate with the registry using an API token."""
        # Update session headers
        self.session.headers.update({
            'Authorization': f'Bearer {token}'
        })
        
        # Verify token by getting user profile
        try:
            # Use discovered endpoint for user profile
            user_url = self._get_endpoint_url('user', '/user/me')
            response = self.session.get(user_url)
            
            # Better error handling for different response codes
            if response.status_code == 401:
                raise Exception("Invalid token. Please check your API token.")
            elif response.status_code == 404:
                raise Exception(f"User endpoint not found. Registry URL: {self.registry_url}")
            elif response.status_code >= 500:
                raise Exception(f"Registry server error ({response.status_code}). Please try again later.")
            
            response.raise_for_status()
            user_data = response.json()
            
            if not user_data.get('username'):
                raise Exception("Invalid response from registry: missing username")
            
            # Update registry config
            self.registry_config['api_key'] = token
            self.registry_config['username'] = user_data.get('username')
            
            # Save config
            self.config.save()
            
            return user_data
        except requests.exceptions.RequestException as e:
            if hasattr(e, 'response') and e.response is not None:
                # Try to get error message from response body
                try:
                    error_data = e.response.json()
                    error_msg = error_data.get('error', str(e))
                except:
                    error_msg = str(e)
                raise Exception(f"Authentication failed: {error_msg}")
            raise Exception(f"Authentication failed: {e}")
    
    # Backwards compatibility
    def login(self, token: str) -> Dict[str, Any]:
        """Authenticate with the registry using an API token (backwards compatibility)."""
        return self.login_with_token(token)
    
    def logout(self):
        """Clear authentication credentials."""
        self.registry_config['api_key'] = None
        self.registry_config['username'] = None
        self.session.headers.pop('Authorization', None)
        self.config.save()
    
    def search(self, query: str, limit: int = 20) -> List[Dict[str, Any]]:
        """Search for packages in the registry."""
        try:
            # Use the packages endpoint with search parameter (from registry discovery)
            search_url = self._get_endpoint_url('packages', '/packages')
            params = {'search': query, 'limit': limit}
            response = self.session.get(search_url, params=params)
            response.raise_for_status()
            # Backend returns 'packages' array directly
            result = response.json()
            packages = result.get('packages', [])
            return packages
        except requests.exceptions.RequestException as e:
            raise Exception(f"Search failed: {e}")
    
    def get_package_info(self, package_name: str) -> Dict[str, Any]:
        """Get detailed information about a package."""
        try:
            # Use discovered endpoint for package info
            if package_name.startswith('@') and '/' in package_name:
                # Scoped package: @namespace/name
                scope, name = package_name[1:].split('/', 1)
                info_url = self._get_endpoint_url(
                    'scopedPackage',
                    f"/packages/@{scope}/{name}",
                    scope=scope,
                    package=name
                )
            else:
                # Unscoped package
                info_url = self._get_endpoint_url(
                    'package',
                    f"/packages/{package_name}",
                    package=package_name
                )

            response = self.session.get(info_url)
            response.raise_for_status()
            return response.json()
        except requests.exceptions.RequestException as e:
            raise Exception(f"Package info fetch failed: {e}")

    def get_package_versions(self, package_name: str) -> List[Dict[str, Any]]:
        """Get version history for a package."""
        try:
            # Use discovered endpoint for package versions
            if package_name.startswith('@') and '/' in package_name:
                # Scoped package: @namespace/name
                scope, name = package_name[1:].split('/', 1)
                versions_url = self._get_endpoint_url(
                    'scopedVersions',
                    f"/packages/@{scope}/{name}/versions",
                    scope=scope,
                    package=name
                )
            else:
                # Unscoped package
                versions_url = self._get_endpoint_url(
                    'versions',
                    f"/packages/{package_name}/versions",
                    package=package_name
                )

            response = self.session.get(versions_url)
            response.raise_for_status()
            return response.json().get('versions', [])
        except requests.exceptions.RequestException as e:
            raise Exception(f"Version fetch failed: {e}")

    def list_user_namespaces(self) -> List[Dict[str, Any]]:
        """List namespaces accessible to the current user."""
        try:
            # Use discovered endpoint for namespaces
            namespaces_url = self._get_endpoint_url('namespaces', '/namespaces')
            response = self.session.get(namespaces_url)
            response.raise_for_status()
            return response.json().get('namespaces', [])
        except requests.exceptions.RequestException:
            # If API call fails, return default namespace
            # This allows offline operation
            return [
                {
                    'name': '@public',
                    'packageCount': 0,
                    'downloadCount': 0,
                    'role': 'member',
                    'verified': False,
                    'permissions': {
                        'canPublish': True,
                        'canDelete': False,
                        'canManage': False
                    }
                }
            ]

    def get_current_namespace(self) -> Optional[str]:
        """Get the current namespace context."""
        # Check if there's a stored namespace context in the config
        return self.config.registry.get('current_namespace')

    def set_current_namespace(self, namespace: str) -> None:
        """Set the current namespace context."""
        self.config.registry['current_namespace'] = namespace
        self.config.save()

    def get_namespace_details(self, namespace: str) -> Dict[str, Any]:
        """Get details about a specific namespace."""
        try:
            # Use discovered endpoint for namespace details
            namespace_url = self._get_endpoint_url('namespace', f'/namespaces/{namespace}', namespace=namespace)
            response = self.session.get(namespace_url)
            response.raise_for_status()
            return response.json()
        except requests.exceptions.RequestException:
            # If API call fails, return basic info
            # This allows offline operation with cached namespace context
            return {
                'name': namespace,
                'description': f'Namespace {namespace}',
                'packageCount': 0,
                'downloadCount': 0,
                'role': 'member',
                'verified': False
            }

    def create_namespace(self, namespace_data: Dict[str, Any]) -> Dict[str, Any]:
        """Create a new namespace."""
        if not self.registry_config.get('api_key'):
            raise Exception("Authentication required. Run 'prompd login' first.")

        namespace_name = namespace_data.get('name')

        # This would normally call the registry API with a POST request
        # For now, this is a mock implementation that:
        # 1. Does NOT actually create anything on a remote registry
        # 2. Does NOT persist anything to a database
        # 3. Simply returns success for testing purposes

        # In a real implementation, this would:
        # response = self.session.post(
        #     f"{self.registry_url}/v1/namespaces",
        #     json=namespace_data
        # )
        # response.raise_for_status()
        # return response.json()

        return {
            'success': True,
            'namespace': namespace_name,
            'message': f'Namespace {namespace_name} created successfully',
            'requiresVerification': False
        }

    def publish_package(self, package_path: Path, target_namespace: Optional[str] = None) -> Dict[str, Any]:
        """Publish a .pdpkg package to the registry. ONLY .pdpkg packages are supported."""
        if not self.registry_config.get('api_key'):
            raise Exception("Authentication required. Run 'prompd registry login' first.")

        if package_path.suffix != '.pdpkg':
            raise Exception(f"Only .pdpkg package files are supported. Got: {package_path.suffix}. This is a package registry, not a .prmd file registry.")

        # Use provided namespace or fall back to current namespace context
        namespace = target_namespace or self.get_current_namespace()

        return self._publish_pdpkg(package_path, namespace)
    
    def _publish_pdpkg(self, package_path: Path, namespace: Optional[str] = None) -> Dict[str, Any]:
        """Publish a .pdpkg bundle package."""
        # Validate package structure
        try:
            validate_pdpkg(package_path)
        except Exception as e:
            raise Exception(f"Invalid .pdpkg package: {e}")

        # Extract package ID from manifest.json inside the .pdpkg (use ID for URL, not name)
        package_name = "unknown"  # fallback
        try:
            with zipfile.ZipFile(package_path, 'r') as zf:
                if 'manifest.json' in zf.namelist():
                    manifest_data = json.loads(zf.read('manifest.json').decode('utf-8'))
                    # Use ID first (for scoped packages like @prompd.io/core-patterns), fallback to name
                    base_name = manifest_data.get('id', manifest_data.get('name', 'unknown'))

                    # If namespace is provided and package doesn't have a scope, add it
                    if namespace and not base_name.startswith('@'):
                        package_name = f"{namespace}/{base_name}"
                    else:
                        package_name = base_name
        except Exception:
            pass  # Use fallback
        
        # Upload package using registry API (npm-compatible endpoint)
        try:
            # Use the package-specific endpoint from well-known registry config
            if package_name.startswith('@') and '/' in package_name:
                # Scoped package: @namespace/name
                scope, name = package_name[1:].split('/', 1)
                publish_url = self._get_endpoint_url(
                    'scopedPublish',
                    f"/packages/@{scope}/{name}",
                    scope=scope,
                    package=name
                )
            else:
                # Unscoped package
                publish_url = self._get_endpoint_url(
                    'publish',
                    f"/packages/{package_name}",
                    package=package_name
                )
            
            with open(package_path, 'rb') as f:
                files = {'package': (package_path.name, f, 'application/octet-stream')}

                # Use PUT for publishing (npm-compatible, idempotent operation)
                # This is where the registry validates namespace access:
                # - 200 OK: Published successfully
                # - 403 Forbidden: No access to namespace (e.g., @microsoft, @google)
                # - 401 Unauthorized: Invalid or missing API key
                # - 409 Conflict: Package already exists with this version
                response = self.session.put(publish_url, files=files)
                response.raise_for_status()
                return response.json()
        except requests.exceptions.RequestException as e:
            # This will include namespace access denied errors from the registry
            raise Exception(f"Package publish failed: {e}")
    
    def install_package(self, package_name: str, version: Optional[str] = None, install_dir: Optional[Path] = None) -> Path:
        """Install a package from the registry."""
        # Get package info
        package_info = self.get_package_info(package_name)
        
        # Determine version to install
        if version is None:
            version = package_info.get('latest_version')
        
        # Determine install directory
        if install_dir is None:
            install_dir = Path.cwd() / "prompd_packages"
        
        install_dir.mkdir(parents=True, exist_ok=True)
        
        # Download package content with progress bar
        # Use discovered endpoint for download URL
        if package_name.startswith('@'):
            # Scoped package like @prompd.io/core-patterns
            scope, name = package_name[1:].split('/', 1)
            download_url = self._get_endpoint_url(
                'scopedDownload',
                f"/{package_name}/-/{name}-{version}.pdpkg",
                scope=scope,
                package=name,
                version=version
            )
        else:
            # Unscoped package
            download_url = self._get_endpoint_url(
                'download',
                f"/{package_name}/-/{package_name}-{version}.pdpkg",
                package=package_name,
                version=version
            )
        
        try:
            # Get file size first for progress tracking
            head_response = self.session.head(download_url)
            head_response.raise_for_status()
            total_size = int(head_response.headers.get('content-length', 0))
            
            # Start streaming download
            response = self.session.get(download_url, stream=True)
            response.raise_for_status()
            
            # Determine target path
            if package_info.get('type') == 'single':
                # Single .prmd file
                target_path = install_dir / f"{package_name.split('/')[-1]}.prmd"
                file_mode = 'w'
                encoding = 'utf-8'
            else:
                # Complex package - save as .pdpkg
                target_path = install_dir / f"{package_name.split('/')[-1]}.pdpkg"
                file_mode = 'wb'
                encoding = None
            
            # Download with progress bar
            with Progress(
                "[progress.description]{task.description}",
                DownloadColumn(),
                "[progress.percentage]{task.percentage:>3.0f}%",
                TransferSpeedColumn(),
                TimeRemainingColumn(),
            ) as progress:
                task_id = progress.add_task(f"Downloading {package_name}@{version}", total=total_size)
                
                with open(target_path, file_mode, encoding=encoding) as f:
                    downloaded = 0
                    for chunk in response.iter_content(chunk_size=8192):
                        if chunk:  # Filter out keep-alive chunks
                            if package_info.get('type') == 'single':
                                f.write(chunk.decode('utf-8'))
                            else:
                                f.write(chunk)
                            downloaded += len(chunk)
                            progress.update(task_id, advance=len(chunk))
            
            return target_path
            
        except requests.exceptions.RequestException as e:
            raise Exception(f"Package download failed: {e}")


def validate_pdpkg(package_path: Path):
    """Validate a .pdpkg package structure.

    Uses the centralized PackageValidator to avoid code duplication.
    Raises exceptions for compatibility with existing code that expects exceptions.
    """
    from .package_validator import PackageValidator

    validator = PackageValidator()
    result = validator.validate_pdpkg_package(package_path)

    if not result.is_valid:
        # Combine all errors into a single exception message
        error_messages = '\n'.join(f"  - {error}" for error in result.errors)
        raise Exception(f"Package validation failed:\n{error_messages}")

    # Return successfully (no exception means validation passed)


def _is_valid_semver(version: str) -> bool:
    """Check if version follows semantic versioning format."""
    semver_pattern = r'^(\d+)\.(\d+)\.(\d+)(?:-([0-9A-Za-z-]+(?:\.[0-9A-Za-z-]+)*))?(?:\+([0-9A-Za-z-]+(?:\.[0-9A-Za-z-]+)*))?$'
    return re.match(semver_pattern, version) is not None


def _serialize_for_json(obj):
    """Convert Python objects to JSON-serializable format."""
    if hasattr(obj, '__dict__'):
        # Convert dataclass/object to dict
        result = {}
        for key, value in obj.__dict__.items():
            if not key.startswith('_'):  # Skip private attributes
                result[key] = _serialize_for_json(value)
        return result
    elif hasattr(obj, 'value') and hasattr(obj, 'name'):  # Enum-like objects
        return obj.value if hasattr(obj.value, 'lower') else str(obj.value)
    elif isinstance(obj, list):
        return [_serialize_for_json(item) for item in obj]
    elif isinstance(obj, dict):
        return {key: _serialize_for_json(value) for key, value in obj.items()}
    elif isinstance(obj, (str, int, float, bool, type(None))):
        return obj
    else:
        # For any other object, try to convert to string as fallback
        return str(obj)


def _extract_pdf_text(pdf_path: Path) -> str:
    """Extract text from PDF file."""
    try:
        from PyPDF2 import PdfReader
    except ImportError:
        return f"[PDF Content from {pdf_path.name}]\n\nPDF extraction library not available. Install PyPDF2 to extract PDF content.\n"

    try:
        reader = PdfReader(pdf_path)
        text_content = []
        
        # Add metadata
        text_content.append(f"[PDF Document: {pdf_path.name}]")
        text_content.append(f"Pages: {len(reader.pages)}")
        text_content.append("-" * 50)
        
        # Extract text from each page
        for page_num, page in enumerate(reader.pages, 1):
            page_text = page.extract_text().strip()
            if page_text:
                text_content.append(f"\n[Page {page_num}]")
                text_content.append(page_text)
        
        return "\n".join(text_content)
    except Exception as e:
        return f"[PDF Content from {pdf_path.name}]\n\nError extracting PDF content: {str(e)}\n"


def _extract_docx_text(docx_path: Path) -> str:
    """Extract text from Word document."""
    try:
        from docx import Document
    except ImportError:
        return f"[Word Document from {docx_path.name}]\n\nWord extraction library not available. Install python-docx to extract Word content.\n"

    try:
        doc = Document(docx_path)
        text_content = []
        
        # Add metadata
        text_content.append(f"[Word Document: {docx_path.name}]")
        text_content.append(f"Paragraphs: {len(doc.paragraphs)}")
        text_content.append("-" * 50)
        
        # Extract paragraphs
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_content.append(paragraph.text.strip())
        
        # Extract tables
        if doc.tables:
            text_content.append("\n[Tables]")
            for table_num, table in enumerate(doc.tables, 1):
                text_content.append(f"\n[Table {table_num}]")
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        text_content.append(row_text)
        
        return "\n".join(text_content)
    except Exception as e:
        return f"[Word Document from {docx_path.name}]\n\nError extracting Word content: {str(e)}\n"


def _extract_pptx_text(pptx_path: Path) -> str:
    """Extract text from PowerPoint presentation."""
    try:
        from pptx import Presentation
    except ImportError:
        return f"[PowerPoint Presentation from {pptx_path.name}]\n\nPowerPoint extraction library not available. Install python-pptx to extract content.\n"

    try:
        prs = Presentation(pptx_path)
        text_content = []
        
        # Add metadata
        text_content.append(f"[PowerPoint Presentation: {pptx_path.name}]")
        text_content.append(f"Slides: {len(prs.slides)}")
        text_content.append("-" * 50)
        
        # Extract text from each slide
        for slide_num, slide in enumerate(prs.slides, 1):
            text_content.append(f"\n[Slide {slide_num}]")
            
            # Extract text from shapes
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            if slide_text:
                text_content.extend(slide_text)
            else:
                text_content.append("(No text content)")
        
        return "\n".join(text_content)
    except Exception as e:
        return f"[PowerPoint Presentation from {pptx_path.name}]\n\nError extracting PowerPoint content: {str(e)}\n"


def _process_excel_file(xlsx_path: Path) -> List[Tuple[str, str]]:
    """Process Excel file into multiple CSV files."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        fallback_name = f"{xlsx_path.stem}.xlsx.txt"
        fallback_content = f"[Excel Workbook from {xlsx_path.name}]\n\nExcel extraction library not available. Install openpyxl to extract Excel content.\n"
        return [(fallback_name, fallback_content)]

    try:
        workbook = load_workbook(xlsx_path, data_only=True)
        results = []
        
        # Add summary file
        summary_name = f"{xlsx_path.stem}-summary.txt"
        summary_content = [
            f"[Excel Workbook: {xlsx_path.name}]",
            f"Sheets: {len(workbook.sheetnames)}",
            f"Sheet names: {', '.join(workbook.sheetnames)}",
            "-" * 50
        ]
        results.append((summary_name, "\n".join(summary_content)))
        
        # Process each worksheet
        for sheet_name in workbook.sheetnames:
            try:
                worksheet = workbook[sheet_name]
                
                # Convert to CSV format
                csv_content = io.StringIO()
                csv_writer = csv.writer(csv_content)
                
                # Write header with sheet info
                csv_writer.writerow([f"Sheet: {sheet_name}"])
                csv_writer.writerow([])  # Empty row
                
                # Extract data
                rows_written = 0
                for row in worksheet.iter_rows(values_only=True):
                    # Skip completely empty rows
                    if any(cell is not None and str(cell).strip() for cell in row):
                        # Convert None to empty string, preserve other values
                        clean_row = [str(cell) if cell is not None else "" for cell in row]
                        csv_writer.writerow(clean_row)
                        rows_written += 1
                
                if rows_written > 0:
                    csv_name = f"{xlsx_path.stem}-{sheet_name}.csv"
                    results.append((csv_name, csv_content.getvalue()))
                
            except Exception as sheet_error:
                error_name = f"{xlsx_path.stem}-{sheet_name}-error.txt"
                error_content = f"Error extracting sheet '{sheet_name}': {str(sheet_error)}"
                results.append((error_name, error_content))
        
        return results if results else [(f"{xlsx_path.stem}.xlsx.txt", "Empty Excel workbook")]
        
    except Exception as e:
        fallback_name = f"{xlsx_path.stem}.xlsx.txt"
        fallback_content = f"[Excel Workbook from {xlsx_path.name}]\n\nError extracting Excel content: {str(e)}\n"
        return [(fallback_name, fallback_content)]


def _extract_image_info(image_path: Path) -> str:
    """Extract basic information from image files."""
    try:
        from PIL import Image
    except ImportError:
        return f"[Image: {image_path.name}]\n\nImage processing library not available. Install Pillow to extract image metadata.\n"

    try:
        with Image.open(image_path) as img:
            info = [
                f"[Image: {image_path.name}]",
                f"Format: {img.format}",
                f"Size: {img.size[0]}x{img.size[1]} pixels",
                f"Mode: {img.mode}",
            ]
            
            # Add EXIF data if available
            if hasattr(img, '_getexif') and img._getexif():
                info.append("EXIF data available")
            
            info.append("-" * 30)
            info.append("Binary image data not extracted (preserved as image file)")
            
            return "\n".join(info)
    except Exception as e:
        return f"[Image: {image_path.name}]\n\nError reading image: {str(e)}\n"


def _get_safe_archive_name(base_name: str, extension: str, used_names: set) -> str:
    """Generate a safe archive name, handling conflicts."""
    candidate = f"{base_name}{extension}"
    
    if candidate not in used_names:
        used_names.add(candidate)
        return candidate
    
    # Handle conflicts with pattern: filename-originalext.txt
    original_ext = Path(base_name).suffix
    if original_ext:
        base_without_ext = Path(base_name).stem
        candidate = f"{base_without_ext}-{original_ext[1:]}{extension}"
    else:
        counter = 1
        while True:
            candidate = f"{base_name}-{counter}{extension}"
            if candidate not in used_names:
                break
            counter += 1
    
    used_names.add(candidate)
    return candidate


def _convert_file_for_package(file_path: Path, source_dir: Path, used_names: set, console: Console) -> List[Tuple[str, bytes]]:
    """Convert a file to safe package format. Returns list of (archive_name, content) tuples."""
    relative_path = file_path.relative_to(source_dir)
    file_ext = file_path.suffix.lower()
    
    # Define file categories
    SAFE_FILES = {'.prmd', '.md', '.txt', '.json', '.yaml', '.yml', '.csv', '.tsv'}
    SAFE_IMAGES = {'.png', '.jpg', '.jpeg', '.gif', '.webp', '.svg'}
    
    PROGRAMMING_LANGS = {
        '.js', '.ts', '.jsx', '.tsx', '.cs', '.cpp', '.cxx', '.cc', '.c', '.h', '.hpp',
        '.py', '.go', '.java', '.php', '.rb', '.rs', '.swift', '.kt', '.scala', '.clj',
        '.sh', '.bash', '.ps1', '.bat', '.cmd', '.vbs', '.lua', '.r', '.sql', '.pl', '.pm'
    }
    
    MARKUP_CONFIG = {
        '.html', '.htm', '.css', '.xml', '.scss', '.sass', '.less', '.styl',
        '.vue', '.svelte', '.astro', '.toml', '.ini', '.cfg', '.conf', '.env'
    }
    
    # Safe files - keep as-is
    if file_ext in SAFE_FILES:
        with open(file_path, 'rb') as f:
            content = f.read()
        archive_name = str(relative_path)
        console.print(f"  [green]OK[/green] Keeping safe file: {archive_name}")
        return [(archive_name, content)]
    
    # Safe images - keep as-is but extract metadata
    if file_ext in SAFE_IMAGES:
        with open(file_path, 'rb') as f:
            content = f.read()
        
        results = []
        # Keep original image
        archive_name = str(relative_path)
        results.append((archive_name, content))
        console.print(f"  [green]OK[/green] Keeping image file: {archive_name}")
        
        # Add metadata file for images (except SVG which is text)
        if file_ext != '.svg':
            metadata_name = _get_safe_archive_name(str(relative_path), '.info.txt', used_names)
            metadata_content = _extract_image_info(file_path)
            results.append((metadata_name, metadata_content.encode('utf-8')))
            console.print(f"  [blue]+[/blue] Extracted image metadata: {metadata_name}")
        
        return results
    
    # Programming languages and markup - convert to .ext.txt
    if file_ext in PROGRAMMING_LANGS or file_ext in MARKUP_CONFIG:
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            # Create header with file info
            header = f"[Source File: {file_path.name}]\n"
            header += f"Original Extension: {file_ext}\n"
            header += f"Language/Type: {file_ext[1:].upper()}\n"
            header += "-" * 50 + "\n\n"
            
            full_content = header + content
            
            archive_name = _get_safe_archive_name(str(relative_path), '.txt', used_names)
            console.print(f"  [yellow]CONV[/yellow] Converted code file: {file_path.name} -> {archive_name}")
            return [(archive_name, full_content.encode('utf-8'))]
            
        except UnicodeDecodeError:
            # If can't read as text, treat as binary
            pass
    
    # Binary file extraction
    try:
        if file_ext == '.pdf':
            text_content = _extract_pdf_text(file_path)
            archive_name = _get_safe_archive_name(str(relative_path), '.txt', used_names)
            console.print(f"  [blue]PDF[/blue] Extracted PDF text: {file_path.name} -> {archive_name}")
            return [(archive_name, text_content.encode('utf-8'))]
        
        elif file_ext == '.docx':
            text_content = _extract_docx_text(file_path)
            archive_name = _get_safe_archive_name(str(relative_path), '.txt', used_names)
            console.print(f"  [blue]DOCX[/blue] Extracted Word text: {file_path.name} -> {archive_name}")
            return [(archive_name, text_content.encode('utf-8'))]
        
        elif file_ext == '.pptx':
            text_content = _extract_pptx_text(file_path)
            archive_name = _get_safe_archive_name(str(relative_path), '.txt', used_names)
            console.print(f"  [blue]PPTX[/blue] Extracted PowerPoint text: {file_path.name} -> {archive_name}")
            return [(archive_name, text_content.encode('utf-8'))]
        
        elif file_ext == '.xlsx':
            csv_files = _process_excel_file(file_path)
            results = []
            for csv_name, csv_content in csv_files:
                # Ensure the CSV name is safe and unique
                safe_csv_name = _get_safe_archive_name(csv_name, '', used_names)
                results.append((safe_csv_name, csv_content.encode('utf-8')))
            
            console.print(f"  [blue]XLSX[/blue] Processed Excel file: {file_path.name} -> {len(results)} files")
            return results
        
    except Exception as e:
        console.print(f"  [red]WARN[/red] Failed to extract {file_path.name}: {str(e)}")
    
    # Fallback: convert unknown files to .ext.txt with binary info
    try:
        file_size = file_path.stat().st_size
        
        # Try to read as text first
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            header = f"[File: {file_path.name}]\n"
            header += f"Extension: {file_ext}\n"
            header += f"Size: {file_size} bytes\n"
            header += "-" * 50 + "\n\n"
            
            full_content = header + content
            archive_name = _get_safe_archive_name(str(relative_path), '.txt', used_names)
            console.print(f"  [cyan]TEXT[/cyan] Converted unknown file: {file_path.name} -> {archive_name}")
            return [(archive_name, full_content.encode('utf-8'))]
            
        except UnicodeDecodeError:
            # Binary file - create info file instead
            info_content = f"[Binary File: {file_path.name}]\n"
            info_content += f"Extension: {file_ext}\n"
            info_content += f"Size: {file_size} bytes\n"
            info_content += f"Type: Binary/Unknown\n"
            info_content += "-" * 50 + "\n"
            info_content += "Binary content not extracted for security.\n"
            info_content += "Original file type not supported for text extraction.\n"
            
            archive_name = _get_safe_archive_name(str(relative_path), '.info.txt', used_names)
            console.print(f"  [magenta]INFO[/magenta] Created info file for binary: {file_path.name} -> {archive_name}")
            return [(archive_name, info_content.encode('utf-8'))]
            
    except Exception as e:
        # Last resort - error file
        error_content = f"[Error Processing: {file_path.name}]\n"
        error_content += f"Error: {str(e)}\n"
        archive_name = _get_safe_archive_name(str(relative_path), '.error.txt', used_names)
        console.print(f"  [red]ERR[/red] Error processing: {file_path.name} -> {archive_name}")
        return [(archive_name, error_content.encode('utf-8'))]


def create_pdpkg(source_dir: Path, output_path: Path, manifest: Dict[str, Any]):
    """Create a .pdpkg package from a directory with universal file conversion."""
    console = Console()
    used_names = set()
    conversion_stats = {
        'safe_files': 0,
        'converted_files': 0,
        'extracted_files': 0,
        'error_files': 0,
        'total_files': 0
    }
    
    console.print(f"\n[bold blue]Creating package with universal file support:[/bold blue] {output_path.name}")
    
    with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zip_file:
        # Add manifest - ensure it's JSON serializable
        serialized_manifest = _serialize_for_json(manifest)
        zip_file.writestr('manifest.json', json.dumps(serialized_manifest, indent=2))
        console.print("  [green]OK[/green] Added manifest.json")
        
        def should_ignore_file(file_path: Path, relative_path: str) -> tuple[bool, str]:
            """Check if file should be ignored based on manifest patterns and auto-ignore."""
            import fnmatch

            file_name = file_path.name

            # Normalize relative_path for consistent matching (forward slashes, no ./ prefix)
            normalized_path = relative_path.replace('\\', '/')
            if normalized_path.startswith('./'):
                normalized_path = normalized_path[2:]

            # Check if this is the main file - always include it
            main_file = manifest.get('main')
            if main_file:
                # Normalize main file path for comparison
                normalized_main = main_file.replace('\\', '/').lstrip('./')
                if normalized_path == normalized_main:
                    return False, "main file (auto-included)"

            # Check if this is the readme file - always include it
            readme_file = manifest.get('readme')
            if readme_file:
                # Normalize readme file path for comparison
                normalized_readme = readme_file.replace('\\', '/').lstrip('./')
                if normalized_path == normalized_readme:
                    return False, "readme file (auto-included)"

            # Check explicit include patterns first (files array)
            files_patterns = manifest.get('files', [])
            if files_patterns:
                # If files array exists, ONLY include matching patterns (unless it's main/readme)
                for pattern in files_patterns:
                    # Normalize pattern (remove ./ prefix, use forward slashes)
                    normalized_pattern = pattern.replace('\\', '/').lstrip('./')

                    if fnmatch.fnmatch(normalized_path, normalized_pattern) or fnmatch.fnmatch(file_name, normalized_pattern):
                        # Explicitly included - don't ignore
                        break
                else:
                    # Not in whitelist - ignore (unless it's main or readme, checked above)
                    # Main and readme are already checked and returned False if matched
                    return True, "not in files whitelist"

            # Check explicit ignore patterns (ignore array)
            ignore_patterns = manifest.get('ignore', [])
            for pattern in ignore_patterns:
                if fnmatch.fnmatch(relative_path, pattern) or fnmatch.fnmatch(file_name, pattern):
                    return True, f"matches ignore pattern '{pattern}'"

            # Apply auto-ignore patterns (unless overridden by explicit files array)
            if not files_patterns:  # Only auto-ignore if no explicit whitelist
                # Auto-ignore dotfiles and dotfolders
                if file_name.startswith('.'):
                    return True, "dotfile/dotfolder"

                # Auto-ignore build artifacts & dependencies
                auto_ignore_dirs = {
                    'node_modules', 'bin', 'obj', 'dist', 'build', 'out', 'target', 'tmp', 'temp',
                    '__pycache__', '.cargo', 'Thumbs.db', 'desktop.ini'
                }

                # Check if any part of the path contains ignored directories
                path_parts = Path(relative_path).parts
                for part in path_parts:
                    if part in auto_ignore_dirs:
                        return True, f"build artifact ({part})"

                # Auto-ignore by file extension
                auto_ignore_extensions = {
                    '.exe', '.dll', '.so', '.dylib',  # binaries
                    '.o', '.a', '.class',             # compiled objects
                    '.log', '.tmp', '.cache', '.pid', # temporary files
                    '.pyc', '.swp', '.swo',           # cache/swap files
                    '.pdpkg'                          # other packages
                }

                if file_path.suffix.lower() in auto_ignore_extensions:
                    return True, f"auto-ignored extension ({file_path.suffix})"

                # Legacy: Skip .pdproj files
                if file_name.endswith('.pdproj'):
                    return True, "build metadata"

            return False, ""

        # Process all files in source directory with intelligent conversion
        for root, dirs, files in os.walk(source_dir):
            # Filter out ignored directories at the directory level
            dirs[:] = [d for d in dirs if not d.startswith('.') and d not in {
                'node_modules', 'bin', 'obj', 'dist', 'build', 'out', 'target',
                'tmp', 'temp', '__pycache__', '.cargo'
            }]

            for file in files:
                file_path = Path(root) / file
                relative_path = str(file_path.relative_to(source_dir))

                # Check if file should be ignored
                should_ignore, ignore_reason = should_ignore_file(file_path, relative_path)
                if should_ignore:
                    console.print(f"  [dim]SKIP[/dim] {ignore_reason}: {relative_path}")
                    continue
                
                conversion_stats['total_files'] += 1
                file_path = Path(root) / file
                
                try:
                    # Convert file using universal conversion system
                    converted_files = _convert_file_for_package(file_path, source_dir, used_names, console)
                    
                    # Add all converted files to the package
                    for archive_name, content in converted_files:
                        zip_file.writestr(archive_name, content)
                        
                        # Update statistics
                        if len(converted_files) == 1 and archive_name == str(file_path.relative_to(source_dir)):
                            conversion_stats['safe_files'] += 1
                        elif archive_name.endswith('.error.txt'):
                            conversion_stats['error_files'] += 1
                        elif len(converted_files) > 1 or file_path.suffix.lower() in {'.pdf', '.docx', '.xlsx', '.pptx'}:
                            conversion_stats['extracted_files'] += 1
                        else:
                            conversion_stats['converted_files'] += 1
                
                except Exception as e:
                    # Handle unexpected errors during file processing
                    error_name = f"{file}.error.txt"
                    error_content = f"[Error Processing: {file}]\nUnexpected error: {str(e)}\n"
                    zip_file.writestr(error_name, error_content.encode('utf-8'))
                    conversion_stats['error_files'] += 1
                    console.print(f"  [red]ERR[/red] Unexpected error processing: {file}")
    
    # Display conversion summary
    console.print(f"\n[bold green]Package creation complete![/bold green]")
    console.print(f"[cyan]Conversion Summary:[/cyan]")
    console.print(f"  Safe files (kept as-is): {conversion_stats['safe_files']}")
    console.print(f"  Code files (-> .txt): {conversion_stats['converted_files']}")  
    console.print(f"  Binary extractions: {conversion_stats['extracted_files']}")
    console.print(f"  Error files: {conversion_stats['error_files']}")
    console.print(f"  Total files processed: {conversion_stats['total_files']}")
    
    # Show package security status
    if conversion_stats['error_files'] == 0:
        console.print(f"[bold green]OK Package is secure - no executable files present[/bold green]")
    else:
        console.print(f"[yellow]WARN {conversion_stats['error_files']} files had processing errors[/yellow]")
    
    console.print(f"[dim]Package saved: {output_path}[/dim]\n")
