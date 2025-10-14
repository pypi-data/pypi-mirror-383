"""
PowerPoint 文件解析器
支持 .pptx 和 .ppt 格式
"""

import asyncio
from pathlib import Path
from typing import List, Dict, Any
from pptx import Presentation

from ..base import BaseParser, ParseResult, FileType


class PPTParser(BaseParser):
    """PowerPoint 文件解析器"""
    
    @property
    def supported_extensions(self) -> List[str]:
        return [".pptx", ".ppt"]
    
    @property
    def file_type(self) -> FileType:
        return FileType.POWERPOINT
    
    async def parse(self, file_path: Path) -> ParseResult:
        """
        解析 PowerPoint 文件
        
        Args:
            file_path: PowerPoint 文件路径
            
        Returns:
            ParseResult: 解析结果
        """
        try:
            if file_path.suffix.lower() == ".pptx":
                text_content = await self._parse_pptx(file_path)
            else:  # .ppt
                text_content = await self._parse_ppt(file_path)
            
            # 获取幻灯片信息
            slide_info = await self._get_slide_info(file_path)
            
            metadata = self._get_file_metadata(file_path)
            metadata.update({
                "ppt_format": file_path.suffix.lower(),
                "slide_count": len(slide_info),
                "slides": slide_info
            })
            
            return ParseResult(
                filename=Path(file_path).name,
                file_type=self.file_type,
                text=text_content.strip(),
                metadata=metadata,
                success=True
            )
            
        except Exception as e:
            self.logger.error(f"PowerPoint 文件解析失败: {str(e)}")
            raise
    
    async def _parse_pptx(self, file_path: Path) -> str:
        """解析 .pptx 文件"""
        def _extract():
            try:
                prs = Presentation(file_path)
                text_parts = []
                
                for slide_num, slide in enumerate(prs.slides, 1):
                    slide_text = []
                    
                    # 提取幻灯片标题
                    title = ""
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            if not title:  # 第一个有文本的形状作为标题
                                title = shape.text.strip()
                            slide_text.append(shape.text.strip())
                    
                    if slide_text:
                        text_parts.append(f"=== 幻灯片 {slide_num} ===")
                        if title:
                            text_parts.append(f"标题: {title}")
                        text_parts.append("\n".join(slide_text))
                        text_parts.append("")  # 空行分隔
                
                return "\n".join(text_parts)
                
            except Exception as e:
                self.logger.error(f"解析 .pptx 文件失败: {str(e)}")
                raise
        
        return await asyncio.get_event_loop().run_in_executor(None, _extract)
    
    async def _parse_ppt(self, file_path: Path) -> str:
        """解析 .ppt 文件（需要额外工具）"""
        def _extract():
            try:
                # .ppt 格式需要额外的库或工具
                # 这里返回提示信息
                return f"[提示] 无法直接解析 .ppt 文件 {file_path.name}，请转换为 .pptx 格式或安装 python-pptx 的 .ppt 支持"
            except Exception as e:
                self.logger.error(f"解析 .ppt 文件失败: {str(e)}")
                return f"[错误] 解析 .ppt 文件失败: {str(e)}"
        
        return await asyncio.get_event_loop().run_in_executor(None, _extract)
    
    async def _get_slide_info(self, file_path: Path) -> List[Dict[str, Any]]:
        """获取幻灯片信息"""
        def _get_info():
            try:
                if file_path.suffix.lower() != ".pptx":
                    return []
                
                prs = Presentation(file_path)
                slide_info = []
                
                for slide_num, slide in enumerate(prs.slides, 1):
                    shapes_count = len(slide.shapes)
                    text_shapes = 0
                    title = ""
                    
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text_shapes += 1
                            if not title:  # 第一个有文本的形状作为标题
                                title = shape.text.strip()
                    
                    slide_info.append({
                        "slide_number": slide_num,
                        "shapes_count": shapes_count,
                        "text_shapes_count": text_shapes,
                        "title": title
                    })
                
                return slide_info
                
            except Exception as e:
                self.logger.error(f"获取幻灯片信息失败: {str(e)}")
                return []
        
        return await asyncio.get_event_loop().run_in_executor(None, _get_info)
