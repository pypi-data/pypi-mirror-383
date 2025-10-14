from typing import List

from pptx import Presentation
from pptx.shapes.group import GroupShape

from pfip.base.constant import TFileExt, TAtomItem
from pfip.base.parser.base import Parser
from pfip.base.parser.parser_model import ParseResult, TitleNode, TitleAtomItem, AtomItem
from pfip.ppt.ppt_ele_handler import PPTTableHandler, PPTNoteHandler, \
    PPTTitleHandler, PPTSlideHandler, PPTShapeHandler, PPTParagraphHandler
from pfip.ppt.ppt_result_repair import PPTRePair


class PPTParser(Parser):
    """
        支持提取PPT的标题,表格及备注信息等

        暂不支持的特性:
        1. 不处理ppt中的图片
        2. 不支持多表头处理
    """
    title_ele_handler: PPTSlideHandler = PPTTitleHandler()
    node_ele_handler: PPTSlideHandler = PPTNoteHandler()
    shape_ele_handlers: List[PPTShapeHandler] = [PPTParagraphHandler(), PPTTableHandler()]
    ppt_repair: PPTRePair = PPTRePair()
    """pdf元素处理器集合"""

    def support(self, file_ext: str) -> bool:
        return file_ext.lower() == TFileExt.PPTX

    def build_titles(self, title_items: List[TitleAtomItem]) -> List[TitleNode]:
        titles = [item.create_and_connect() for item in title_items]
        self.fill_titles(titles)
        return titles

    def process_shape(self, shape, page_num: int) -> List[AtomItem]:
        rtn = []
        for shape_handler in self.shape_ele_handlers:
            if shape_handler.support(shape):
                rtn.extend(shape_handler.run(shape, page_num))
        return rtn

    def __call__(self, file_path: str, **kwargs) -> ParseResult:
        items = [TitleAtomItem.from_file(file_path)]
        prs = Presentation(file_path)
        for idx, slide in enumerate(prs.slides):
            page_num = idx + 1
            if self.title_ele_handler.support(slide):
                items.extend(self.title_ele_handler.run(slide, page_num))
            for shape in slide.shapes:
                if isinstance(shape, GroupShape):
                    for son_shape in shape.shapes:
                        items.extend(self.process_shape(son_shape, page_num))
                else:
                    items.extend(self.process_shape(shape, page_num))
            if self.node_ele_handler.support(slide):
                items.extend(self.node_ele_handler.run(slide, page_num))
        items = self.ppt_repair.run(items)
        title_items = [item for item in items if item.item_type == TAtomItem.TITLE]
        titles = self.build_titles(title_items)
        return ParseResult(
            items=items,
            titles=titles
        )
