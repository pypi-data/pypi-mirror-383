from abc import abstractmethod, ABC
from typing import List

import pandas as pd
from pandas import DataFrame
from pptx.shapes.autoshape import Shape
from pptx.shapes.graphfrm import GraphicFrame
from pptx.slide import Slide
from pydantic import BaseModel

from pfip.base.errors import UnImplementException
from pfip.base.parser.parser_model import AtomItem, TextAtomItem, TableAtomItem, TitleAtomItem


class PPTSlideHandler(ABC, BaseModel):
    @abstractmethod
    def support(self, slide: Slide) -> bool:
        raise UnImplementException()

    @abstractmethod
    def run(self, slide: Slide, page_num: int) -> List[AtomItem]:
        raise UnImplementException()


class PPTShapeHandler(ABC, BaseModel):
    @abstractmethod
    def support(self, shape: Shape) -> bool:
        raise UnImplementException()

    @abstractmethod
    def run(self, shape: Shape, page_num: int) -> List[AtomItem]:
        raise UnImplementException()


class PPTParagraphHandler(PPTShapeHandler):
    def support(self, shape: Shape) -> bool:
        return shape.has_text_frame

    def run(self, shape: Shape, page_num: int) -> List[TextAtomItem]:
        rtn = []
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                rtn.append(TextAtomItem(content=run.text, start_page=page_num, end_page=page_num))
        return rtn


class PPTTableHandler(PPTShapeHandler):
    def support(self, shape: Shape) -> bool:
        return shape.has_table

    @staticmethod
    def to_data_frame(shape: GraphicFrame) -> DataFrame:
        table = shape.table
        data = []
        columns = []
        for i, row in enumerate(table.rows):
            if i == 0:
                columns.extend([cell.text for cell in row.cells])
            else:
                data.append([cell.text for cell in row.cells])
        return pd.DataFrame(data, columns=columns)

    def run(self, shape: GraphicFrame, page_num: int) -> List[TableAtomItem]:
        df = self.to_data_frame(shape)
        table_item = TableAtomItem.simple(df=df, page_num=page_num)
        return [table_item]


class PPTTitleHandler(PPTSlideHandler):
    def support(self, slide: Slide) -> bool:
        return slide.shapes.title is not None

    def run(self, slide: Slide, page_num: int) -> List[TitleAtomItem]:
        slide_title = slide.shapes.title.text
        title_item = TitleAtomItem(
            content=slide_title,
            start_page=page_num,
            end_page=page_num,
            title_level=1,
        )
        return [title_item]


class PPTNoteHandler(PPTSlideHandler):
    def support(self, slide: Slide) -> bool:
        return slide.has_notes_slide

    def run(self, slide: Slide, page_num: int) -> List[TextAtomItem]:
        text = slide.notes_slide.notes_text_frame.text
        node_item = TextAtomItem(content=text, start_page=page_num, end_page=page_num)
        return [node_item]
